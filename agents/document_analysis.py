"""
Document Analysis Agent for RFP and File Processing
"""

import os
import tempfile
import zipfile
import json
from datetime import datetime, timezone
from typing import Dict, List, Any, Optional
from werkzeug.utils import secure_filename

from .base import BaseAgent, AgentResult

class DocumentAnalysisAgent(BaseAgent):
    """Agent specialized in document analysis, text extraction, and file processing"""
    
    def __init__(self):
        super().__init__(
            name="Document Analysis Agent",
            description="Extracts text from files, analyzes documents, and processes RFP uploads",
            capabilities=["text_extraction", "file_processing", "document_analysis", "rfp_processing"],
            persona_icon="📄",
            persona_color="#ff4d4d"
        )
        self.agent_type = "ai"
    
    def execute(self, task: str, context: Dict[str, Any] = None) -> Dict[str, Any]:
        """Execute document analysis task"""
        self.log_usage()
        
        if context is None:
            context = {}
        
        try:
            if task == "extract_text":
                return self._extract_text(context)
            elif task == "process_rfp_upload":
                return self._process_rfp_upload(context)
            elif task == "analyze_document":
                return self._analyze_document(context)
            else:
                return AgentResult(
                    success=False,
                    error=f"Unknown task: {task}",
                    agent_id=self.id
                ).to_dict()
        except Exception as e:
            return AgentResult(
                success=False,
                error=str(e),
                agent_id=self.id,
                function_used=task
            ).to_dict()
    
    def get_available_functions(self) -> List[Dict[str, Any]]:
        """Get available functions for this agent"""
        return [
            {
                "name": "extract_text",
                "description": "Extract text from a file (PDF, DOCX, TXT)",
                "parameters": {
                    "filepath": "Path to the file to extract text from"
                }
            },
            {
                "name": "process_rfp_upload",
                "description": "Process an uploaded RFP file and extract requirements",
                "parameters": {
                    "file_data": "File data or file path",
                    "filename": "Name of the uploaded file"
                }
            },
            {
                "name": "analyze_document",
                "description": "Analyze document content and extract key information",
                "parameters": {
                    "text": "Text content to analyze",
                    "analysis_type": "Type of analysis to perform"
                }
            }
        ]
    
    def _extract_text_from_file(self, filepath: str) -> str:
        """Extract text from various file formats"""
        try:
            ext = os.path.splitext(filepath)[1].lower()
            
            if ext == '.txt':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            
            elif ext == '.pdf':
                try:
                    # Try PyMuPDF first (better extraction)
                    import fitz  # PyMuPDF
                    doc = fitz.open(filepath)
                    text = ""
                    for page in doc:
                        text += page.get_text() + "\n"
                    doc.close()
                    return text
                except ImportError:
                    try:
                        # Fall back to PyPDF2
                        import PyPDF2
                        with open(filepath, 'rb') as f:
                            reader = PyPDF2.PdfReader(f)
                            text = ""
                            for page in reader.pages:
                                text += page.extract_text() + "\n"
                        return text
                    except ImportError:
                        return "Error: No PDF library installed (PyMuPDF or PyPDF2 required)"
            
            elif ext == '.docx':
                try:
                    import docx
                    doc = docx.Document(filepath)
                    text = ""
                    for paragraph in doc.paragraphs:
                        text += paragraph.text + "\n"
                    # Also extract text from tables
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                text += cell.text + " "
                            text += "\n"
                    return text
                except ImportError:
                    return "Error: python-docx not installed"
            
            elif ext == '.doc':
                try:
                    # Try textract first
                    import textract
                    return textract.process(filepath).decode('utf-8', errors='ignore')
                except ImportError:
                    try:
                        # Try antiword as fallback
                        import subprocess
                        result = subprocess.run(['antiword', filepath], 
                                              capture_output=True, text=True, errors='ignore')
                        return result.stdout if result.returncode == 0 else "Error: Could not process .doc file"
                    except (ImportError, FileNotFoundError):
                        return "Error: textract or antiword required for .doc files"
            
            elif ext in ['.ppt', '.pptx']:
                try:
                    from pptx import Presentation
                    prs = Presentation(filepath)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return text
                except ImportError:
                    return "Error: python-pptx not installed"
            
            elif ext == '.rtf':
                try:
                    from striprtf.striprtf import rtf_to_text
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        rtf = f.read()
                    return rtf_to_text(rtf)
                except ImportError:
                    # Basic fallback - just read as text (won't properly handle RTF formatting)
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        return f.read()
            
            elif ext == '.odt':
                try:
                    import odf.opendocument
                    import odf.text
                    doc = odf.opendocument.load(filepath)
                    paragraphs = doc.getElementsByType(odf.text.P)
                    return '\n'.join([p.plainText() for p in paragraphs])
                except ImportError:
                    try:
                        import textract
                        return textract.process(filepath).decode('utf-8', errors='ignore')
                    except ImportError:
                        return "Error: No ODT library installed"
            
            elif ext in ['.html', '.htm']:
                try:
                    from bs4 import BeautifulSoup
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        soup = BeautifulSoup(f.read(), 'html.parser')
                    return soup.get_text(separator='\n')
                except ImportError:
                    # Raw HTML as fallback
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        return f.read()
            
            elif ext in ['.md', '.markdown']:
                try:
                    import markdown
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        md_text = f.read()
                    html = markdown.markdown(md_text)
                    try:
                        from bs4 import BeautifulSoup
                        return BeautifulSoup(html, 'html.parser').get_text()
                    except ImportError:
                        return md_text  # Return raw markdown as fallback
                except ImportError:
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        return f.read()  # Return raw markdown
            
            elif ext in ['.xlsx', '.xls']:
                try:
                    import pandas as pd
                    df = pd.read_excel(filepath)
                    return df.to_string()
                except ImportError:
                    return "Error: pandas and openpyxl required for Excel files"
            
            elif ext == '.csv':
                try:
                    import pandas as pd
                    df = pd.read_csv(filepath)
                    return df.to_string()
                except ImportError:
                    # Basic fallback
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        return f.read()
            
            else:
                return f"Unsupported file type: {ext}"
                
        except Exception as e:
            return f"Error extracting text: {str(e)}"
    
    def _extract_text(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Extract text from a file"""
        filepath = context.get("filepath")
        if not filepath:
            return AgentResult(
                success=False,
                error="No filepath provided",
                agent_id=self.id,
                function_used="extract_text"
            ).to_dict()
        
        if not os.path.exists(filepath):
            return AgentResult(
                success=False,
                error=f"File not found: {filepath}",
                agent_id=self.id,
                function_used="extract_text"
            ).to_dict()
        
        extracted_text = self._extract_text_from_file(filepath)
        
        return AgentResult(
            success=True,
            data={
                "text": extracted_text,
                "filepath": filepath,
                "extracted_at": datetime.now(timezone.utc).isoformat()
            },
            agent_id=self.id,
            function_used="extract_text"
        ).to_dict()
    
    def _process_rfp_upload(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Process an uploaded RFP file"""
        file_data = context.get("file_data")
        filename = context.get("filename", "unknown.txt")
        
        if not file_data:
            return AgentResult(
                success=False,
                error="No file data provided",
                agent_id=self.id,
                function_used="process_rfp_upload"
            ).to_dict()
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as temp_file:
            temp_file.write(file_data)
            temp_filepath = temp_file.name
        
        try:
            # Extract text from the uploaded file
            extracted_text = self._extract_text_from_file(temp_filepath)
            
            # Create processing log
            processing_log = {
                "timestamp": datetime.now(timezone.utc).isoformat(),
                "filename": filename,
                "file_size": len(file_data),
                "extraction_success": not extracted_text.startswith("Error:"),
                "text_length": len(extracted_text) if not extracted_text.startswith("Error:") else 0
            }
            
            return AgentResult(
                success=True,
                data={
                    "text": extracted_text,
                    "filename": filename,
                    "processing_log": processing_log
                },
                agent_id=self.id,
                function_used="process_rfp_upload"
            ).to_dict()
            
        finally:
            # Clean up temp file
            if os.path.exists(temp_filepath):
                os.unlink(temp_filepath)
    
    def _analyze_document(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze document content"""
        text = context.get("text", "")
        analysis_type = context.get("analysis_type", "general")
        
        if not text:
            return AgentResult(
                success=False,
                error="No text provided for analysis",
                agent_id=self.id,
                function_used="analyze_document"
            ).to_dict()
        
        # Basic analysis
        analysis_result = {
            "word_count": len(text.split()),
            "character_count": len(text),
            "line_count": len(text.split('\n')),
            "analysis_type": analysis_type,
            "analyzed_at": datetime.now(timezone.utc).isoformat()
        }
        
        # Add specific analysis based on type
        if analysis_type == "rfp":
            analysis_result.update(self._analyze_rfp_content(text))
        
        return AgentResult(
            success=True,
            data=analysis_result,
            agent_id=self.id,
            function_used="analyze_document"
        ).to_dict()
    
    def _analyze_rfp_content(self, text: str) -> Dict[str, Any]:
        """Analyze RFP-specific content"""
        # Simple keyword-based analysis
        keywords = {
            "requirements": ["requirement", "must", "shall", "should", "need"],
            "deadlines": ["deadline", "due", "submit", "delivery"],
            "budget": ["budget", "cost", "price", "payment", "$"],
            "qualifications": ["qualification", "experience", "certification", "skill"]
        }
        
        analysis = {}
        for category, words in keywords.items():
            count = sum(text.lower().count(word) for word in words)
            analysis[f"{category}_mentions"] = count
        
        return analysis
